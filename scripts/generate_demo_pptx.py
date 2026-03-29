#!/usr/bin/env python3
"""
Generate Resolvly pitch deck (PowerPoint) — aligned with demo_script.md and deck visuals:
  • Cover: Track 1, tagline, Indiana → every American, problem stats
  • Who we’re building for: kitchen-table quote + Primary / Secondary / Tertiary
  • Technical execution → live demo handoff
  • Ethics + three questions
  • What’s next: Phases 2–3 pillars (Polish & Reach · Expand Coverage · Sustainable Impact)

Requires: pip install python-pptx

Usage:
  pip install python-pptx
  python scripts/generate_demo_pptx.py
  open scripts/Resolvly_Pitch_Demo.pptx
"""
from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    raise SystemExit("Install python-pptx: pip install python-pptx") from None


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs: Presentation, title: str, bullets: list[str], font_pt: int = 18) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Resolvly",
        "Track 1 • Biology & Physical Health • Claude Hackathon 2026\n"
        "Navigating health insurance denials\n"
        "A free, AI-powered tool for Indiana patients — and eventually, every American",
    )

    add_bullets(
        prs,
        "The problem (by the numbers)",
        [
            "~1 in 5 insured Americans face a denied claim.",
            "<1% of denials are ever appealed.",
            "99% of patients don’t know they can fight back.",
            "Gap: jargon & codes — not lack of caring. Resolvly closes it with plain English + action.",
        ],
        font_pt=17,
    )

    add_bullets(
        prs,
        "Who we’re building for — and why it matters",
        [
            "“The person who opens a denial letter at their kitchen table and has no idea what it means "
            "or what to do next — facing a bill they shouldn’t owe.”",
            "PRIMARY: Uninsured / underinsured Indiana patients with denial letters and no clear next step.",
            "SECONDARY: Small medical practices without dedicated billing teams to chase denials.",
            "TERTIARY: Patient advocates & social workers helping vulnerable populations navigate insurance.",
            "Why it scales: Predictable denial/EOB shapes nationwide + live authoritative lookups (CMS, regs, state DOI).",
        ],
        font_pt=13,
    )

    add_bullets(
        prs,
        "Technical execution · Purposeful AI → Live demo",
        [
            "Working prototype: upload denial + EOB + bill; plan context → extract → enrich (public sources) → "
            "analyze (root cause, deadlines, outputs).",
            "AI used purposefully: unstructured text, classification, drafting + rules/structured extraction — "
            "assistive, not a single black-box score.",
            "Stack: React + FastAPI; session-forward MVP (no long-term PHI warehouse by default).",
            "→ Next: live product walkthrough.",
        ],
        font_pt=14,
    )

    add_bullets(
        prs,
        "Ethics — harms & what we ship today",
        [
            "Data: MVP avoids a claims DB; session-based UI; server processes in memory for the request (v1 no PHI-at-rest story).",
            "UI: Footer disclaimers on analyze, action plan, bill breakdown, appeal drafting, results — not a law firm; not legal/medical advice.",
            "Appeal drafting: “Edit before sending”; recommend attorney/advocate before insurer submission.",
            "Transparency: Code/reg enrichment from public sources; assumptions panel + denial completeness where exposed.",
            "Autonomy: No auto-filing; user reviews and decides. Three Qs: who for, what could go wrong, help vs decide.",
        ],
        font_pt=11,
    )

    add_bullets(
        prs,
        "What’s next — Phases 2–3 (weeks 5–12) · V2 adds persistence (e.g. Supabase)",
        [
            "Polish & reach: plain-English for low digital literacy; Spanish first; mobile camera & clipboard upload; WCAG AA.",
            "Expand coverage: 50-state deadline logic; DOI complaint/review integration; anonymous denial-trends dataset; "
            "provider exports.",
            "Sustainable impact: provider-portal / B2B path; advocacy & policy; advocacy org partnerships; HIPAA BAA & "
            "grants/nonprofit path when data persists.",
            "Impact: help people use appeal rights they already have — access to care & financial health.",
            "Thank you — questions?",
        ],
        font_pt=12,
    )

    out = Path(__file__).resolve().parent / "Resolvly_Pitch_Demo.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
